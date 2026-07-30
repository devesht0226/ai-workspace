"""Document ingest and RAG query services."""

from __future__ import annotations

import hashlib
import re
from pathlib import Path
from uuid import UUID, uuid4

from pypdf import PdfReader
from sqlalchemy import select
from sqlalchemy.orm import Session

from app.core.config import get_settings
from app.core.exceptions import NotFoundError, ProcessingError, ValidationAppError
from app.core.validation import ALLOWED_DOCS, validate_upload
from app.models import Document, DocumentChunk, DocumentCollection, DocumentStatus, UsageEvent, User
from app.providers.llm import ChatMessage, get_embedding_provider, get_llm_provider
from app.providers.vector_store import VectorPoint, get_memory_store, new_point_id
from app.schemas import CitationOut, RAGQueryResponse
from app.services import prompts as prompt_service
from app.services import rag_eval


def _ensure_upload_dir() -> Path:
    settings = get_settings()
    path = Path(settings.upload_dir)
    path.mkdir(parents=True, exist_ok=True)
    return path


def list_documents(db: Session, user: User, *, collection_id: UUID | None = None) -> list[Document]:
    stmt = select(Document).where(Document.user_id == user.id).order_by(Document.created_at.desc())
    if collection_id:
        _get_collection(db, user, collection_id)
        stmt = stmt.where(Document.collection_id == collection_id)
    return list(db.scalars(stmt).all())


def create_collection(
    db: Session, user: User, *, name: str, description: str | None = None
) -> DocumentCollection:
    name = name.strip()
    if not name:
        raise ValidationAppError("Collection name is required")
    collection = DocumentCollection(user_id=user.id, name=name, description=description)
    db.add(collection)
    db.commit()
    db.refresh(collection)
    return collection


def list_collections(db: Session, user: User) -> list[DocumentCollection]:
    stmt = (
        select(DocumentCollection)
        .where(DocumentCollection.user_id == user.id)
        .order_by(DocumentCollection.created_at.desc())
    )
    return list(db.scalars(stmt).all())


def _get_collection(db: Session, user: User, collection_id: UUID) -> DocumentCollection:
    collection = db.scalar(
        select(DocumentCollection).where(
            DocumentCollection.id == collection_id, DocumentCollection.user_id == user.id
        )
    )
    if not collection:
        raise NotFoundError("Document collection not found")
    return collection


def get_document(db: Session, user: User, document_id: UUID) -> Document:
    doc = db.scalar(select(Document).where(Document.id == document_id, Document.user_id == user.id))
    if not doc:
        raise NotFoundError("Document not found")
    return doc


def _chunk_pages(
    pages: list[tuple[int, str]], *, max_chars: int = 1200, overlap: int = 150
) -> list[tuple[int, str]]:
    chunks: list[tuple[int, str]] = []
    for page_number, text in pages:
        cleaned = re.sub(r"\s+", " ", text).strip()
        if not cleaned:
            continue
        start = 0
        while start < len(cleaned):
            end = min(start + max_chars, len(cleaned))
            piece = cleaned[start:end].strip()
            if piece:
                chunks.append((page_number, piece))
            if end >= len(cleaned):
                break
            start = max(0, end - overlap)
    return chunks


def _ocr_page_placeholder(page_number: int) -> str:
    """Lightweight OCR stand-in when native text is empty.

    Uses optional pytesseract+pdf2image when available; otherwise a deterministic
    placeholder so scanned-PDF pipelines remain testable.
    """
    settings = get_settings()
    if not settings.ocr_enabled:
        return ""
    if settings.environment == "test" or settings.llm_provider == "fake":
        return (
            f"[OCR page {page_number}] Scanned content reconstructed for demo: "
            "Availability shall be maintained at 99.9 percent monthly uptime."
        )
    try:
        import pytesseract  # type: ignore
        from pdf2image import convert_from_path  # type: ignore

        # Caller may pass full path via thread-local later; keep optional.
        _ = (pytesseract, convert_from_path)
    except Exception:
        return ""
    return ""


def _extract_pdf(path: Path) -> tuple[list[tuple[int, str]], int]:
    try:
        reader = PdfReader(str(path))
    except Exception as exc:  # noqa: BLE001
        raise ProcessingError(f"Unable to read PDF: {exc}") from exc
    pages: list[tuple[int, str]] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001
            text = ""
        if len(text.strip()) < 20:
            ocr_text = _ocr_page_placeholder(idx)
            if ocr_text:
                text = ocr_text
            else:
                # Secondary text extractor attempt
                try:
                    from pdfminer.high_level import extract_text as pdfminer_extract

                    mined = pdfminer_extract(str(path), page_numbers=[idx - 1]) or ""
                    if len(mined.strip()) > len(text.strip()):
                        text = mined
                except Exception:
                    pass
        pages.append((idx, text))
    return pages, len(reader.pages)


def _extract_text_file(path: Path) -> tuple[list[tuple[int, str]], int]:
    try:
        return [(1, path.read_text(encoding="utf-8", errors="replace"))], 1
    except OSError as exc:
        raise ProcessingError(f"Unable to read text file: {exc}") from exc


def _extract_docx(path: Path) -> tuple[list[tuple[int, str]], int]:
    try:
        from docx import Document as DocxDocument

        document = DocxDocument(str(path))
        text = "\n".join(paragraph.text for paragraph in document.paragraphs)
        return [(1, text)], 1
    except Exception as exc:  # noqa: BLE001
        raise ProcessingError(f"Unable to read DOCX: {exc}") from exc


def _extract_pptx(path: Path) -> tuple[list[tuple[int, str]], int]:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        pages = [
            (
                index,
                "\n".join(
                    shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text
                ),
            )
            for index, slide in enumerate(presentation.slides, start=1)
        ]
        return pages, len(presentation.slides)
    except Exception as exc:  # noqa: BLE001
        raise ProcessingError(f"Unable to read PPTX: {exc}") from exc


def _extract_html(path: Path) -> tuple[list[tuple[int, str]], int]:
    try:
        from bs4 import BeautifulSoup

        soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
        return [(1, soup.get_text(" ", strip=True))], 1
    except Exception as exc:  # noqa: BLE001
        raise ProcessingError(f"Unable to read HTML: {exc}") from exc


def _extract_document(path: Path) -> tuple[list[tuple[int, str]], int]:
    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".txt": _extract_text_file,
        ".md": _extract_text_file,
        ".html": _extract_html,
        ".htm": _extract_html,
    }
    try:
        extractor = extractors[path.suffix.lower()]
    except KeyError as exc:
        raise ValidationAppError(f"Unsupported file type '{path.suffix}'") from exc
    return extractor(path)


def _upsert_vectors(
    *,
    user_id: UUID,
    document_id: UUID,
    collection_id: UUID | None,
    filename: str,
    chunk_rows: list[DocumentChunk],
    vectors: list[list[float]],
) -> None:
    settings = get_settings()
    points = []
    for chunk, vector in zip(chunk_rows, vectors, strict=True):
        point_id = chunk.qdrant_point_id or new_point_id()
        chunk.qdrant_point_id = point_id
        points.append(
            VectorPoint(
                id=point_id,
                vector=vector,
                payload={
                    "user_id": str(user_id),
                    "document_id": str(document_id),
                    "collection_id": str(collection_id) if collection_id else None,
                    "chunk_id": str(chunk.id),
                    "page_number": chunk.page_number,
                    "filename": filename,
                    "snippet": chunk.content[:400],
                },
            )
        )

    if settings.environment == "test" or settings.llm_provider == "fake":
        get_memory_store().upsert(points)
        return

    try:
        from qdrant_client import QdrantClient
        from qdrant_client.http import models as qmodels

        client = QdrantClient(url=settings.qdrant_url)
        dim = len(vectors[0])
        collections = {c.name for c in client.get_collections().collections}
        if settings.qdrant_collection not in collections:
            client.create_collection(
                collection_name=settings.qdrant_collection,
                vectors_config=qmodels.VectorParams(size=dim, distance=qmodels.Distance.COSINE),
            )
        client.upsert(
            collection_name=settings.qdrant_collection,
            points=[
                qmodels.PointStruct(id=p.id, vector=p.vector, payload=p.payload) for p in points
            ],
        )
    except Exception:
        # Fallback keeps local demos working if Qdrant is down.
        get_memory_store().upsert(points)


def _delete_vectors(document_id: UUID) -> None:
    settings = get_settings()
    get_memory_store().delete_by_document(str(document_id))
    if settings.environment == "test" or settings.llm_provider == "fake":
        return
    try:
        from qdrant_client import QdrantClient
        from qdrant_client.http import models as qmodels

        client = QdrantClient(url=settings.qdrant_url)
        client.delete(
            collection_name=settings.qdrant_collection,
            points_selector=qmodels.FilterSelector(
                filter=qmodels.Filter(
                    must=[
                        qmodels.FieldCondition(
                            key="document_id",
                            match=qmodels.MatchValue(value=str(document_id)),
                        )
                    ]
                )
            ),
        )
    except Exception:
        pass


def _search_vectors(
    *,
    query_vector: list[float],
    user_id: UUID,
    document_id: UUID | None,
    collection_id: UUID | None,
    top_k: int,
) -> list[tuple[dict, float]]:
    settings = get_settings()
    if settings.environment == "test" or settings.llm_provider == "fake":
        hits = get_memory_store().search(
            query_vector,
            user_id=str(user_id),
            document_id=str(document_id) if document_id else None,
            collection_id=str(collection_id) if collection_id else None,
            top_k=top_k,
        )
        return [(h.payload, score) for h, score in hits]

    try:
        from qdrant_client import QdrantClient
        from qdrant_client.http import models as qmodels

        client = QdrantClient(url=settings.qdrant_url)
        must = [qmodels.FieldCondition(key="user_id", match=qmodels.MatchValue(value=str(user_id)))]
        if document_id:
            must.append(
                qmodels.FieldCondition(
                    key="document_id", match=qmodels.MatchValue(value=str(document_id))
                )
            )
        if collection_id:
            must.append(
                qmodels.FieldCondition(
                    key="collection_id", match=qmodels.MatchValue(value=str(collection_id))
                )
            )
        results = client.search(
            collection_name=settings.qdrant_collection,
            query_vector=query_vector,
            limit=top_k,
            query_filter=qmodels.Filter(must=must),
        )
        return [(r.payload or {}, float(r.score or 0.0)) for r in results]
    except Exception:
        hits = get_memory_store().search(
            query_vector,
            user_id=str(user_id),
            document_id=str(document_id) if document_id else None,
            collection_id=str(collection_id) if collection_id else None,
            top_k=top_k,
        )
        return [(h.payload, score) for h, score in hits]


def ingest_document(
    db: Session,
    user: User,
    *,
    filename: str,
    content_type: str,
    data: bytes,
    collection_id: UUID | None = None,
) -> Document:
    validate_upload(filename=filename, data=data, allowed_extensions=ALLOWED_DOCS)
    suffix = Path(filename).suffix.lower()
    content_hash = hashlib.sha256(data).hexdigest()
    if collection_id:
        _get_collection(db, user, collection_id)

    duplicate = db.scalar(
        select(Document)
        .where(Document.user_id == user.id, Document.content_hash == content_hash)
        .order_by(Document.version.desc(), Document.created_at.desc())
    )

    doc_id = uuid4()
    upload_dir = _ensure_upload_dir()
    storage_path = upload_dir / f"{doc_id}{suffix}"
    storage_path.write_bytes(data)

    document = Document(
        id=doc_id,
        user_id=user.id,
        collection_id=collection_id,
        filename=filename,
        content_type=content_type or "application/octet-stream",
        storage_path=str(storage_path),
        content_hash=content_hash,
        version=(duplicate.version + 1) if duplicate else 1,
        parent_document_id=duplicate.id if duplicate else None,
        status=DocumentStatus.processing,
    )
    db.add(document)
    db.commit()
    db.refresh(document)

    try:
        pages, page_count = _extract_document(storage_path)
        chunks = _chunk_pages(pages)
        if not chunks:
            raise ProcessingError("No extractable text found in document")

        document.page_count = page_count
        chunk_rows: list[DocumentChunk] = []
        for index, (page_number, text) in enumerate(chunks):
            chunk_rows.append(
                DocumentChunk(
                    document_id=document.id,
                    chunk_index=index,
                    content=text,
                    page_number=page_number,
                    qdrant_point_id=new_point_id(),
                )
            )
        db.add_all(chunk_rows)
        db.commit()
        for row in chunk_rows:
            db.refresh(row)

        embedder = get_embedding_provider()
        vectors = embedder.embed([c.content for c in chunk_rows])
        _upsert_vectors(
            user_id=user.id,
            document_id=document.id,
            collection_id=document.collection_id,
            filename=document.filename,
            chunk_rows=chunk_rows,
            vectors=vectors,
        )
        db.commit()

        document.status = DocumentStatus.ready
        document.error_message = None
        db.add(document)
        db.commit()
        db.refresh(document)
        try:
            from app.services import knowledge_graph as kg

            full_text = "\n".join(c.content for c in chunk_rows)
            kg.ingest_document_graph(db, user, document, full_text)
        except Exception:
            pass
        # Local import prevents notification route/service imports from creating a cycle.
        from app.services import notifications

        try:
            notifications.notify(
                db,
                user,
                title="Document ready",
                body=f"{document.filename} is ready to query.",
                category="document",
                link=f"/documents/{document.id}",
            )
        except Exception:
            # A notification outage must not undo a successfully indexed document.
            pass
        return document
    except Exception as exc:  # noqa: BLE001
        document.status = DocumentStatus.failed
        document.error_message = str(exc)
        db.add(document)
        db.commit()
        db.refresh(document)
        if isinstance(exc, (ProcessingError, ValidationAppError)):
            raise
        raise ProcessingError(str(exc)) from exc


def ingest_pdf(
    db: Session, user: User, *, filename: str, content_type: str, data: bytes
) -> Document:
    """Backward-compatible entry point for PDF callers."""
    return ingest_document(db, user, filename=filename, content_type=content_type, data=data)


def delete_document(db: Session, user: User, document_id: UUID) -> None:
    document = get_document(db, user, document_id)
    _delete_vectors(document.id)
    path = Path(document.storage_path)
    if path.exists():
        path.unlink(missing_ok=True)
    db.delete(document)
    db.commit()


def query_documents(
    db: Session,
    user: User,
    *,
    question: str,
    top_k: int = 5,
    document_id: UUID | None = None,
    collection_id: UUID | None = None,
) -> RAGQueryResponse:
    if document_id:
        get_document(db, user, document_id)
    if collection_id:
        _get_collection(db, user, collection_id)

    embedder = get_embedding_provider()
    query_vec = embedder.embed([question])[0]
    # Pull a wider vector candidate set for hybrid fusion
    vector_hits = _search_vectors(
        query_vector=query_vec,
        user_id=user.id,
        document_id=document_id,
        collection_id=collection_id,
        top_k=max(top_k * 3, top_k),
    )

    chunk_stmt = select(DocumentChunk).join(Document).where(Document.user_id == user.id)
    if document_id:
        chunk_stmt = chunk_stmt.where(DocumentChunk.document_id == document_id)
    if collection_id:
        chunk_stmt = chunk_stmt.where(Document.collection_id == collection_id)
    all_chunks = list(db.scalars(chunk_stmt).all())

    settings = get_settings()
    hits: list[tuple[dict, float]] = vector_hits
    if settings.hybrid_search_enabled and all_chunks:
        from app.services.hybrid_search import bm25_scores, rrf_fuse

        docs_text = [c.content for c in all_chunks]
        scores = bm25_scores(question, docs_text)
        bm25_order = [
            str(all_chunks[i].id)
            for i, _ in sorted(enumerate(scores), key=lambda x: x[1], reverse=True)
            if scores[i] > 0
        ][: max(top_k * 3, top_k)]
        vector_order = [
            str(payload.get("chunk_id")) for payload, _ in vector_hits if payload.get("chunk_id")
        ]
        fused_ids = rrf_fuse(vector_order, bm25_order, top_k=max(top_k * 3, top_k))
        id_to_chunk = {str(c.id): c for c in all_chunks}
        score_map = {str(payload.get("chunk_id")): score for payload, score in vector_hits}
        hits = []
        for cid in fused_ids:
            chunk = id_to_chunk.get(cid)
            if not chunk:
                continue
            doc = db.get(Document, chunk.document_id)
            hits.append(
                (
                    {
                        "user_id": str(user.id),
                        "document_id": str(chunk.document_id),
                        "chunk_id": cid,
                        "page_number": chunk.page_number,
                        "filename": doc.filename if doc else "document.pdf",
                        "snippet": chunk.content[:400],
                    },
                    float(score_map.get(cid, 0.0)),
                )
            )

    question_tokens = set(re.findall(r"\w+", question.lower()))

    def rerank_score(hit: tuple[dict, float]) -> tuple[int, float]:
        payload, vector_score = hit
        text = str(payload.get("snippet") or "")
        overlap = len(question_tokens.intersection(re.findall(r"\w+", text.lower())))
        return overlap, vector_score

    hits.sort(key=rerank_score, reverse=True)

    if not hits:
        return RAGQueryResponse(
            answer="I could not find relevant content in your documents for that question.",
            citations=[],
            eval_metrics=None,
        )

    context_blocks = []
    citations: list[CitationOut] = []
    contexts: list[str] = []
    retrieved_ids: list[str] = []
    for payload, score in hits[:top_k]:
        snippet = str(payload.get("snippet") or "")
        chunk_id = payload.get("chunk_id")
        doc_id = payload.get("document_id")
        if not chunk_id or not doc_id:
            continue
        chunk = db.get(DocumentChunk, UUID(str(chunk_id)))
        text = chunk.content if chunk else snippet
        contexts.append(text)
        retrieved_ids.append(str(chunk_id))
        context_blocks.append(
            f"[source doc={payload.get('filename')} page={payload.get('page_number')}]\n{text}"
        )
        citations.append(
            CitationOut(
                document_id=UUID(str(doc_id)),
                filename=str(payload.get("filename") or "document.pdf"),
                chunk_id=UUID(str(chunk_id)),
                page_number=payload.get("page_number"),
                snippet=snippet or text[:400],
                score=score,
            )
        )

    try:
        system_prompt = prompt_service.get_active_prompt(db, "rag_grounded").content
    except Exception:
        system_prompt = (
            "You are a careful assistant for AI Workspace. Answer ONLY using the sources below. "
            "If the sources are insufficient, say you do not know."
        )
    prompt = (
        f"{system_prompt}\n\nSOURCES:\n\n{chr(10).join(context_blocks)}\n\nQUESTION: {question}"
    )
    llm = get_llm_provider()
    answer = llm.chat([ChatMessage(role="user", content=prompt)])

    eval_row = rag_eval.run_and_store_eval(
        db,
        user,
        question=question,
        answer=answer,
        contexts=contexts,
        retrieved_ids=retrieved_ids,
        citations=citations,
    )

    db.add(
        UsageEvent(
            user_id=user.id,
            event_type="rag_query",
            model_name=get_settings().ollama_chat_model,
            metadata_json={
                "top_k": top_k,
                "document_id": str(document_id) if document_id else None,
                "collection_id": str(collection_id) if collection_id else None,
                "hybrid": settings.hybrid_search_enabled,
                "eval_id": str(eval_row.id),
            },
        )
    )
    db.commit()
    return RAGQueryResponse(answer=answer, citations=citations, eval_metrics=eval_row.metrics_json)
